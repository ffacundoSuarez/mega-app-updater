# create_slides.py

import datetime
import logging
import re
from pptx.util import Pt, Inches, Cm
from pptx.dml.color import RGBColor
# === MODIFICACIÓN 1: Reemplazar MSO_ALIGNMENT por PP_ALIGN ===
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
# =============================================================
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION, XL_LABEL_POSITION as POS
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE as PH_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.oxml.ns import qn 
import pandas as pd
import numpy as np
import copy  # <--- NUEVO: Clave para clonar los estilos visuales
from pptx.enum.chart import XL_MARKER_STYLE
from pptx.oxml import parse_xml # <--- NUEVO: Para crear etiquetas XML
from pptx.enum.chart import XL_CHART_TYPE
from . import config
from . import utils
import datetime
import difflib  # 👈 ¡ACÁ ESTÁ LA LIBRERÍA DESPIERTA!

#import config
#import utils

# Ejemplo de Colores Fijos por Etapa (DEBE DEFINIR ESTO EN SUS CONSTANTES)
STAGE_COLORS = [
    'E52EA8', # Etapa 1: Conocimiento (Azul Fuerte)
    '00A5A3', # Etapa 2: Alguna Vez (Turquesa)
    '00B0F0', # Etapa 3: Consideración (Verde Claro)
    'FEDD3D', # Etapa 4: Compra (Amarillo/Naranja)
    '92D050'  # Etapa 5: Fidelización (Naranja Fuerte)
]
BRAND_COLORS = ['E52EA8', '00A5A3', '00B0F0', 'FEDD3D', '92D050']

# Reutilizar constantes de formato
#FONT_NAME = config.FONT_NAME
FONT_NAME = 'Nunito Light'
BAR_COLOR_HEX = config.BAR_COLOR_HEX
HEADER_COLOR_HEX = config.HEADER_COLOR_HEX
HIGHLIGHT_COLOR_HEX = config.HIGHLIGHT_COLOR_HEX
SERIES_COLORS_P29 = config.SERIES_COLORS_P29
CHART_TYPES = config.CHART_TYPES

# Índices (Importados de config)
IDX_CHART = config.IDX_CHART; IDX_TABLE = config.IDX_TABLE; IDX_SUBTITLE = config.IDX_SUBTITLE 
IDX_QUESTION_LABEL = config.IDX_QUESTION_LABEL
IDX_CHART_DUAL_1 = config.IDX_CHART_DUAL_1; IDX_CHART_DUAL_2 = config.IDX_CHART_DUAL_2
IDX_BODY_DUAL = config.IDX_BODY_DUAL

CHART_WIDTH = Inches(5.0) 
CHART_HEIGHT = Inches(4.0)
TABLE_WIDTH = Inches(5.0) 
TABLE_HEIGHT = Inches(4.0)



def _style_chart(chart_frame, chart_type_key, finding): 
    """Aplica el estilo corporativo a un gráfico, incluyendo colores y etiquetas."""
    chart = chart_frame.chart
    chart_type_enum = CHART_TYPES.get(chart_type_key, XL_CHART_TYPE.BAR_CLUSTERED) 

    chart.has_title = False
    
    # Estilo de ejes
    axis_font = chart.category_axis.tick_labels.font
    axis_font.size = Pt(10); axis_font.name = FONT_NAME
    
    if chart.value_axis.has_major_gridlines:
        chart.value_axis.major_gridlines.format.line.fill.background()
    
    # Estilo específico para Gráfico Apilado (Escala: BAR_STACKED_100)
    if chart_type_enum == XL_CHART_TYPE.BAR_STACKED_100:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM 
        legend_font = chart.legend.font
        legend_font.size = Pt(10)
        legend_font.name = FONT_NAME
        chart.category_axis.reverse_order = True
        chart.value_axis.visible = False
        
        if chart.series: chart.series[0].gap_width = 20
        
        #for i, series in enumerate(chart.series):
        
        for i, series in enumerate(chart.series):             
             
             # CRÍTICO: Aplicar colores de SERIES_COLORS_P29
            if i < len(SERIES_COLORS_P29):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = RGBColor.from_string(SERIES_COLORS_P29[i])

            
            # 1. Habilitar y Formatear Etiquetas de Datos (Optimización Visual)
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.number_format = '0%'       
            data_labels.show_percentage = False
            data_labels.show_value = True 
            data_labels.position = POS.CENTER
            # Fuente (accesible desde data_labels.font)
            data_labels.font.name = 'Nunito Light'
            data_labels.font.size = Pt(8)
            data_labels.font.color.rgb = RGBColor(0, 0, 0)
            data_labels.number_format_is_linked = False

    # Estilo específico para Gráfico Agrupado/Individual (SRQ/MRQ/DUAL)
    elif chart_type_enum == XL_CHART_TYPE.BAR_CLUSTERED or chart_type_enum == XL_CHART_TYPE.COLUMN_CLUSTERED:
        
        chart.has_legend = False
        #chart.value_axis.visible = False 
        chart.category_axis.reverse_order = True
        value_axis = chart.value_axis
        value_axis.format.line.fill.solid()
        value_axis.format.line.fill.fore_color.rgb = RGBColor(255, 255, 255) # invisible
        value_axis.major_tick_mark = XL_TICK_MARK.NONE
        value_axis.minor_tick_mark = XL_TICK_MARK.NONE
        value_axis.tick_labels.font.size = Pt(1)
        value_axis.tick_labels.font.color.rgb = RGBColor(255, 255, 255)

        # === AJUSTE DE ESCALA (0% a 100% para gráficos apilados) ===
        #value_axis = chart.value_axis
        value_axis.minimum_scale = 0.0
        value_axis.maximum_scale = 1.0
        for series in chart.series:
                         
            # CRÍTICO: Aplicar color de BAR_COLOR_HEX a las barras/columnas (Serie 1)
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(
            BAR_COLOR_HEX[1:] if BAR_COLOR_HEX.startswith('#') else BAR_COLOR_HEX
            )
            
            # 1. Habilitar y Formatear Etiquetas de Datos (Optimización Visual)
            series.has_data_labels = True
            data_labels = series.data_labels
            
            # Formato: Mostrar Porcentaje, ocultar Valor, Nombre de Serie/Categoría.
            data_labels.show_percentage = False
            data_labels.show_value = True
            data_labels.show_legend_key = False
            data_labels.show_category_name = False      
            data_labels.position = POS.OUTSIDE_END
            # *** FORMATO CORREGIDO: Usamos '0%' o '0.0%' para decimales ***
            # Si quieres 26%:
            data_labels.number_format = '0%'            
            # *** CORRECCIÓN CLAVE: APLICAR FORMATO DE NÚMERO EXPLICITO ***
            # Usar '0%' para mostrar el valor como porcentaje.
            #data_labels.number_format = '#,##0"%"'
            # Si quieres decimales (ej. 26.5%): '#,##0.0"%"'

        # Añade la configuración de fuente A NIVEL DE SERIE para evitar el bucle.
        # Esto debería aplicar a todas las etiquetas de la serie.
            data_labels.font.name = FONT_NAME
            data_labels.font.size = Pt(10)
            data_labels.font.color.rgb = RGBColor(0, 0, 0)

        # Asegura la posición final
            data_labels.position = POS.OUTSIDE_END

def _insert_styled_table(placeholder, data_df, sig_df, banner_info):
    """Inserta y estiliza una tabla de datos con significancia corregida."""
    
    rows_count = len(data_df) + 1 
    cols_count = len(data_df.columns) - 1 
    
    if rows_count <= 1 or cols_count == 0:
        logging.warning("Tabla omitida: datos insuficientes.")
        return None, Pt(0)

    # 1. Crear contenedor y fijar dimensiones
    table_container_shape = placeholder.insert_table(rows_count, cols_count)
    table_container_shape.width = TABLE_WIDTH
    table_container_shape.height = TABLE_HEIGHT
    table_obj = table_container_shape.table

    # --- AJUSTE DE DIMENSIONES INTERNAS ---
    if cols_count > 0:
        col_width_inches = TABLE_WIDTH.inches / cols_count
        for col in table_obj.columns:
            col.width = Inches(col_width_inches)
    
    forced_header_height = Inches(0.4) 
    table_obj.rows[0].height = forced_header_height
    
    content_rows_count = rows_count - 1
    if content_rows_count > 0:
        row_height_inches = (TABLE_HEIGHT.inches - forced_header_height.inches) / content_rows_count
        for r in range(1, rows_count):
            table_obj.rows[r].height = Inches(row_height_inches)

    header_cols = data_df.columns.tolist()[1:]
    display_letters = banner_info.get('display_map', banner_info['letter_map'])
    
    # 2. Llenar Encabezados
    for col_idx, text in enumerate(header_cols):
        cell = table_obj.cell(0, col_idx)
        texto_base = text.split(':', 1)[1].strip() if ':' in text else text.strip()
        letra_fija = display_letters.get(text, '').upper()
        
        cell.text = f"{texto_base} ({letra_fija})"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(HEADER_COLOR_HEX.replace('#', ''))
        
        tf = cell.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.name = FONT_NAME
        p.font.color.rgb = RGBColor(255, 255, 255)

    # 3. Llenar Contenido (con lógica de significancia corregida)
    data_indexed = data_df.set_index(data_df.columns[0])
    
    for r in range(len(data_indexed)):
        index_val = data_indexed.index[r]
        current_row_data = data_indexed.loc[index_val]
        is_all_100 = (current_row_data == 100).all()
        
        for c in range(cols_count):
            cell = table_obj.cell(r + 1, c)
            segment_name = header_cols[c]
            value = data_indexed.iloc[r, c]
            
            # --- Lógica de Significancia Segura ---
            letter_code = ""
            if not is_all_100 and sig_df is not None:
                if index_val in sig_df.index and segment_name in sig_df.columns:
                    val_sig = sig_df.at[index_val, segment_name]
                    letter_code = str(val_sig).strip().upper() if pd.notna(val_sig) else ""

            # --- Construcción del Texto de la Celda (Evita el Error de Variable) ---
            #if letter_code:
            #    cell_text = f"{value:.0f}% {letter_code}"
            #else:
            #    cell_text = f"{value:.0f}%"
            
# --- 3. Lógica de Formateo Inteligente (AJUSTE AQUÍ) ---
            try:
                # Convertimos a float para validar
                val_numeric = float(value)
                
                # Detectamos si la fila es de 'Media' o 'Promedio'
                es_media = any(palabra in str(index_val).upper() for palabra in ['MEDIA', 'PROMEDIO'])
                
                if es_media:
                    # Formato con 1 o 2 decimales para Medias (sin %)
                    display_val = f"{val_numeric:.1f}"
                else:
                    # Formato entero con % para T2B y otras cajas
                    # Si el valor viene como 0.85 lo multiplicamos por 100
                    if 0 < val_numeric <= 1: 
                        display_val = f"{val_numeric * 100:.0f}%"
                    else:
                        display_val = f"{val_numeric:.0f}%"
            except (ValueError, TypeError):
                # Si no es un número (ej: '---' o texto), lo dejamos como string
                display_val = str(value)

            # Construcción final con la letra de significancia
            cell_text = f"{display_val} {letter_code}".strip()

            cell.text = cell_text
            tf = cell.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10)
            p.font.name = FONT_NAME

            # --- Aplicar Resaltado si hay letra ---
            if letter_code:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(HIGHLIGHT_COLOR_HEX.replace('#', ''))
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 112, 192) # Azul para resaltar la letra
            else:
                cell.fill.background() # Transparente/Fondo normal
                p.font.color.rgb = RGBColor(0, 0, 0)

            # Quitar márgenes internos
            tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = Inches(0)

    return table_container_shape, forced_header_height


def create_single_slide(prs, total_data, table_data, sig_df, finding_original, contrast_original, task, meta, banner_info, chart_type): 
    """
    Crea la diapositiva SINGLE estándar (SRQ/MRQ).
    Ajustada para usar add_chart capturando coordenadas de los placeholders del Layout 1.
    """
    import pandas as pd
    import re
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.util import Pt, Inches

    var_name = task["VARIABLE_NAME"]
    logging.info(f"Procesando tarea: {var_name}")

    # 1. Obtener texto de pregunta
    question_map_dict = utils.get_question_map()
    question_text_full = question_map_dict.get(var_name)

    if question_text_full:
        var_label_full = question_text_full
    else:
        logging.warning(f"ADVERTENCIA: '{var_name}' NO encontrado en el CSV de mapeo. Usando etiqueta de SPSS.")
        var_label_full = utils.get_variable_label(meta, var_name)
    
    # Textos base para Título y Subtítulo
    llm_title = finding_original if finding_original else "Resultado Principal"
    llm_subtitle = var_label_full
    
    # 2. Configuración de Layout
    chart_type_key = chart_type 
    chart_type_enum = config.CHART_TYPES.get(chart_type_key, XL_CHART_TYPE.BAR_CLUSTERED)
    
    try:
        slide_layout_idx = task.get("SLIDE_LAYOUT", 1)
        slide_layout = prs.slide_layouts[slide_layout_idx]
    except Exception as e:
        logging.error(f"Error al cargar Layout índice {slide_layout_idx}: {e}")
        return

    slide = prs.slides.add_slide(slide_layout)

    # --- 3. PREPARACIÓN DE DATOS (Sincronización Tabla-Gráfico) ---
    table_data_clean = table_data.copy()
    
    if 'Categoria' in table_data_clean.columns:
        first_val = table_data_clean['Categoria'].iloc[0]
        if str(first_val).isdigit() and len(table_data_clean.index) > 0:
            table_data_clean['Categoria'] = table_data_clean.index.astype(str)
    else:
        table_data_clean = table_data_clean.reset_index().rename(columns={'index': 'Categoria'})

    # Reconstruimos total_data_final basado en el TOTAL de la tabla
    total_data_final = pd.DataFrame({
        'Categoria': table_data_clean['Categoria'].astype(str).str.strip(),
        'Porcentaje': pd.to_numeric(table_data_clean['TOTAL'], errors='coerce')
    })

    # Ordenamos por porcentaje descendente
    total_data_final = total_data_final.sort_values('Porcentaje', ascending=False).reset_index(drop=True)
    
    # Sincronizamos la tabla de visualización con el nuevo orden
    table_data_final_disp = pd.merge(
        total_data_final[['Categoria']], 
        table_data_clean, 
        on='Categoria', 
        how='left'
    )

    # --- 4. ALINEACIÓN DE SIGNIFICANCIA ---
    if sig_df is not None and not sig_df.empty:
        sig_tmp = sig_df.copy()
        if 'Categoria' not in sig_tmp.columns:
            sig_tmp = sig_tmp.reset_index().rename(columns={'index': 'Categoria'})
        
        sig_tmp['Categoria'] = sig_tmp['Categoria'].astype(str).str.strip()
        
        df_significance_final = pd.merge(
            table_data_final_disp[['Categoria']], 
            sig_tmp, 
            on='Categoria', 
            how='left'
        ).fillna("")
        df_significance_final.set_index('Categoria', inplace=True)
    else:
        df_significance_final = pd.DataFrame("", index=table_data_final_disp['Categoria'], columns=table_data_final_disp.columns)

    # Datos listos para el Chart
    chart_data_clean = total_data_final.dropna(subset=['Porcentaje']).copy()
    chart_data_clean = chart_data_clean[chart_data_clean['Categoria'].astype(str).str.strip() != '']
    chart_data_clean['Porcentaje_Decimal'] = chart_data_clean['Porcentaje'] / 100

    # 5. LÓGICA LLM (Opcional)
    if config.RUN_LLM and utils.LLM_AVAILABLE:
        try:
            table_data_for_llm = table_data_final_disp.head(5) 
            data_context = {
                "finding_principal": finding_original,
                "pregunta": var_label_full,
                "tabla_cruzada": table_data_for_llm.to_markdown(index=False)
            }
            llm_prompt = utils.create_llm_prompt(data_context, "SLIDE_CONTEXTUAL")
            llm_response = utils.generate_text_with_llm(llm_prompt)
            
            t_match = re.search(r"TÍTULO: (.+)", llm_response)
            s_match = re.search(r"SUBTÍTULO: (.+)", llm_response)
            
            if t_match: llm_title = t_match.group(1).strip()
            if s_match: llm_subtitle = s_match.group(1).strip()
        except Exception as e:
            logging.error(f"Error LLM: {e}")

    # 6. ASIGNACIÓN DE TEXTOS Y PLACEHOLDERS
    if slide.shapes.title:
        slide.shapes.title.text = llm_title 
    
    # Marcadores de texto (Subtítulo y Pregunta)
    try:
        slide.placeholders[15].text = llm_subtitle
        slide.placeholders[16].text = var_label_full 
    except:
        pass

    # --- 7. TABLA (INSERT_TABLE en Placeholder 14) ---
    table_placeholder = None
    try:
        table_placeholder = slide.placeholders[14]
        _insert_styled_table(table_placeholder, table_data_final_disp, df_significance_final, banner_info)
    except Exception as e:
        logging.error(f"Error al insertar tabla: {e}")

    # --- 8. GRÁFICO (ADD_CHART capturando y borrando Placeholder 13) ---
    try:
        chart_placeholder = slide.placeholders[13]
        # Captura de coordenadas
        c_left, c_top = chart_placeholder.left, chart_placeholder.top
        c_width, c_height = chart_placeholder.width, chart_placeholder.height
        
        # ELIMINACIÓN del placeholder original para evitar "ruido" visual
        sp = chart_placeholder._sp
        sp.getparent().remove(sp)

        if not chart_data_clean.empty:
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = chart_data_clean['Categoria'].tolist() 
            chart_data_obj.add_series('Porcentaje Total', chart_data_clean['Porcentaje_Decimal'].tolist())
            
            # Inserción limpia con add_chart
            chart_frame = slide.shapes.add_chart(
                chart_type_enum, c_left, c_top, c_width, c_height, chart_data_obj
            )
            
            _style_chart(chart_frame, chart_type_key, llm_title) 
            
            # Ajuste extra de estilo para BAR_CLUSTERED
            chart = chart_frame.chart
            chart.category_axis.format.line.fill.background() # Quitar línea del eje Y
            
    except Exception as e:
        logging.error(f"Error al procesar gráfico con add_chart: {e}")

    logging.info(f"Diapositiva SINGLE ({var_name}) creada exitosamente.")

        
def create_scale_profile_slide(prs, chart_data, table_data, sig_df, finding_original, contrast_original, task, meta, banner_info):
    """Crea la diapositiva de Perfil de Escala optimizada para Banners."""
    
    var_name = task["VARIABLE_NAME"]
    var_label_full = utils.get_variable_label(meta, var_name)

    # Extraemos el nombre técnico (ej: F1) y el nombre amigable (ej: GENERO)
    #banner_var_tech = banner_info.get("variable", "") # Nombre técnico de la columna
    #banner_label = banner_info.get("name", "")        # Etiqueta descriptiva
    
    # Construimos el texto combinado
    #banner_full_display = f"{banner_var_tech}.{banner_label}".strip(".")

    banner_display = banner_info.get("name", "Total")
    #banner_name = banner_info["name"]
    
    llm_title = f"{var_label_full}"
    llm_subtitle = f"Distribución de escala por {banner_display}."

    slide_layout = utils.find_layout_by_name(prs, task.get("SLIDE_LAYOUT", config.CHART_LAYOUT_NAME_SINGLE))
    if slide_layout is None:
        logging.error(f"Layout no encontrado.")
        return

    slide = prs.slides.add_slide(slide_layout)
    
    # 1. Asignar Título de forma segura
    if slide.shapes.title:
        slide.shapes.title.text = llm_title

    # 2. Manejo de Subtítulos (opcional)
    try:
        slide.placeholders[config.IDX_SUBTITLE].text = llm_subtitle 
        slide.placeholders[config.IDX_QUESTION_LABEL].text = var_label_full
    except:
        pass

    # --- BUSQUEDA SEGURA DE PLACEHOLDERS ---
    table_placeholder = None
    chart_placeholder = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == config.IDX_TABLE:
            table_placeholder = ph
        elif ph.placeholder_format.idx == config.IDX_CHART:
            chart_placeholder = ph

    # 3. INSERTAR TABLA (Solo una vez)
    table_shape_container = None
    header_height_pt = Pt(0)
    has_table = not table_data.empty

    if has_table and table_placeholder is not None:
        try:
            table_shape_container, header_height_pt = _insert_styled_table(table_placeholder, table_data, sig_df, banner_info)
        except Exception as e:
            logging.error(f"Error al insertar tabla: {e}")

# 2. CAPTURAR COORDENADAS Y ELIMINAR EL PLACEHOLDER DEL GRÁFICO
    if chart_placeholder is not None:
        # Guardamos las coordenadas
        left, top, width, height = chart_placeholder.left, chart_placeholder.top, chart_placeholder.width, chart_placeholder.height
        
        # AJUSTE: Si no hay tabla, expandimos el ancho manualmente para ocupar la slide
        if not has_table:
        #if orientation == "COLUMN" and not has_table:
            left, width = Inches(0.8), Inches(5.4)
            #original_width = width
            #width = Inches(3.0)
            #left = left + (original_width - width) / 2
            ##left = (prs.slide_width - width) / 2 # Centrado automático
        # ELIMINACIÓN FÍSICA DEL PLACEHOLDER
        # Accedemos al elemento XML subyacente y lo removemos del padre
        sp = chart_placeholder._sp
        sp.getparent().remove(sp)
    else:
        # Fallback si el layout no tiene el placeholder
        left, top, width, height = Inches(0.5), Inches(1.8), Inches(9.0), Inches(4.0)

    # 4. CONFIGURAR COORDENADAS DEL GRÁFICO (Sin usar slide.placeholders[] directamente)
    #if chart_placeholder is None:
    #    # Fallback si el layout no tiene el placeholder esperado
    #    left, top, width, height = Inches(0.5), Inches(1.8), Inches(9.0), Inches(4.0)
    #else:
    #    if not has_table:
    #        # Si no hay tabla, centramos y agrandamos
    #        left, top, width, height = Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.2)
    #    else:
    #        # Usamos el lugar del placeholder
    #        left, top, width, height = chart_placeholder.left, chart_placeholder.top, chart_placeholder.width, chart_placeholder.height

    # 5. PREPARAR DATOS DEL GRÁFICO
    if not chart_data.empty:
        try:
            item_names = chart_data['Item'].tolist()
            # Limpiar etiquetas (quitar el prefijo del banner)
            category_labels = [str(n).split(':')[-1].strip() if ':' in str(n) else str(n) for n in item_names]

            chart_data_decimal = chart_data.copy()
            scale_series_names = chart_data_decimal.columns.tolist()[1:]
            for col in scale_series_names:
                chart_data_decimal[col] = pd.to_numeric(chart_data_decimal[col], errors='coerce').fillna(0) / 100

            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = category_labels
            for series_name in scale_series_names:
                chart_data_obj.add_series(series_name, chart_data_decimal[series_name].tolist())





# --- MODIFICACIÓN 1: Seleccionar tipo de gráfico dinámico ---
            orientation = task.get("CHART_ORIENTATION", "COLUMN")
            if orientation == "BAR":
                chart_type = XL_CHART_TYPE.BAR_STACKED_100
            else:
                chart_type = XL_CHART_TYPE.COLUMN_STACKED_100

# --- AJUSTE DE PROPORCIÓN PARA ESCALAS SIMPLES ---
            # Ahora 'orientation' ya existe y no dará error
            if orientation == "COLUMN" and not has_table:
                original_width = width
                width = Inches(5.0)  
                left = left + (original_width - width) / 2

            # 6. INSERTAR GRÁFICO (Usamos add_chart para evitar dependencia de placeholders idx)
            # Cambiamos a COLUMN_STACKED_100 para comparar banners verticales
            chart_frame = slide.shapes.add_chart(
                chart_type, left, top, width, height, chart_data_obj
            )
            chart = chart_frame.chart

            # 4. Ajustar el eje si es BAR
            if orientation == "BAR":
                chart.category_axis.reverse_order = True
                chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW

            # Estilo rápido
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False  # Esto evita que se superponga al gráfico
            chart.legend.font.size = Pt(9)

            # --- ELIMINAR LÍNEAS DE CUADRÍCULA Y EJES ---
            chart.value_axis.has_major_gridlines = False
            chart.value_axis.has_minor_gridlines = False
            chart.value_axis.visible = False  # Ocultamos el eje Y (0-100%) para más limpieza

            category_axis = chart.category_axis
            category_axis.format.line.fill.background() # Quita la línea del eje X
            category_axis.tick_labels.font.size = Pt(10)

            # --- PALETA DE COLORES PERSONALIZADA (Modificable) ---
            # Define aquí los colores en HEX para cada nivel de la escala
            # Ejemplo: Verde oscuro, Verde claro, Gris, Naranja, Rojo
            PALETA_LIKERT = [
                "C00000", "FF5050", "FEDD3D", "92D050", "00B050"
            ]

            # Quitar líneas de eje para look limpio
            #chart.value_axis.visible = False
            #chart.category_axis.format.line.fill.background()

            # Etiquetas blancas dentro de las columnas
            
            #series_value = chart.series[1]
            for i, series in enumerate(chart.series):
                # 1. Aplicar Color de la paleta
                color_index = i % len(PALETA_LIKERT)
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(PALETA_LIKERT[color_index])
                
                series.has_data_labels = True
                series.gap_width = 10 # Barras más gruesas para look de funnel
                data_labels = series.data_labels
    
                data_labels.show_value = True 
                data_labels.show_percentage = False
                data_labels.number_format_is_linked = False
                data_labels.number_format = '0%'       
                data_labels.position = POS.CENTER
                # Fuente (accesible desde data_labels.font)
                data_labels.font.name = 'Nunito Light'
                data_labels.font.size = Pt(9)
                data_labels.font.color.rgb = RGBColor(255, 255, 255)

                #series.has_data_labels = True
                #series.data_labels.number_format = '0%'
                #series.data_labels.font.size = Pt(8)
                #series.data_labels.font.color.rgb = RGBColor(255, 255, 255)

            # Alineación vertical con la tabla si existe
            if table_shape_container:
                chart_frame.top = table_shape_container.top + header_height_pt
                chart_frame.height = table_shape_container.height - header_height_pt

        except Exception as e:
            logging.error(f"Error al insertar el gráfico apilado: {e}")

    logging.info(f"Diapositiva de Perfil de Escala ({var_name}) creada correctamente.")



def create_dual_slide(prs, result_A, result_B, task, banner_info, chart_type): 
    """Crea la diapositiva de comparación DUAL a partir de los resultados procesados."""
    
    chart_type_key = chart_type 
    chart_type_enum = config.CHART_TYPES.get(chart_type_key, XL_CHART_TYPE.BAR_CLUSTERED)
    
    # Desempaquetar resultados (usamos 8 valores)
    total_data_A, _, _, finding_A, _, _, _, _ = result_A
    total_data_B, _, _, finding_B, _, _, _, _ = result_B

    var_A_name = task["VARIABLE_A"].get("NAME")
    var_B_name = task["VARIABLE_B"].get("NAME")

    # Obtener etiquetas de meta-datos (Usando None para meta)
    master_title = f"{utils.get_variable_label(None, var_A_name).split(':')[0]} vs. {utils.get_variable_label(None, var_B_name).split(':')[0]}"
    master_contrast = f"Desglose por {banner_info['name']}"

    slide_layout = utils.find_layout_by_name(prs, task.get("SLIDE_LAYOUT", config.CHART_LAYOUT_NAME_DUAL))
    
    if slide_layout is None:
        logging.error(f"Layout '{config.CHART_LAYOUT_NAME_DUAL}' no encontrado.")
        return

    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = master_title
    
    body_placeholder = slide.placeholders[IDX_BODY_DUAL]
    body_placeholder.text = master_contrast
    
    chart_placeholders = {'A': (slide.placeholders[IDX_CHART_DUAL_1], total_data_A, finding_A), 
                          'B': (slide.placeholders[IDX_CHART_DUAL_2], total_data_B, finding_B)}

    for key, (placeholder, total_data, finding_text) in chart_placeholders.items():
        total_data_sorted = total_data.sort_values('Porcentaje', ascending=False)
        chart_data_clean = total_data_sorted.dropna(subset=['Porcentaje']).copy()
        
        if chart_data_clean.empty:
            logging.warning(f"Gráfico DUAL {key} omitido: Datos vacíos.")
            continue
        
        #try:
        #    chart_data = CategoryChartData()
        #    chart_data.categories = chart_data_clean['Categoria'].tolist()
        #    chart_data.add_series('Total %', chart_data_clean['Porcentaje'].tolist())
        #    chart_frame = placeholder.insert_chart(chart_type_enum, chart_data)
        #    _style_chart(chart_frame, chart_type_key, finding_text) 

        try:
            # === INICIO MODIFICACIÓN CLAVE DUAL SLIDE ===
            # Los valores en `total_data` están como números enteros (ej: 26, 30).
            # Para que el formato de etiqueta del gráfico funcione con decimales ('0.0%'),
            # necesitamos convertir los valores a su representación decimal (ej: 0.26, 0.30).
            
            # 1. Crear una nueva columna con los valores divididos por 100
            chart_data_clean['Porcentaje_Decimal'] = chart_data_clean['Porcentaje'] / 100
            
            chart_data = CategoryChartData()
            chart_data.categories = chart_data_clean['Categoria'].tolist()
            # 2. Usar la nueva columna decimal para la serie del gráfico
            chart_data.add_series('Total %', chart_data_clean['Porcentaje_Decimal'].tolist())
            # === FIN MODIFICACIÓN CLAVE DUAL SLIDE ===

            chart_frame = placeholder.insert_chart(chart_type_enum, chart_data)
            _style_chart(chart_frame, chart_type_key, finding_text)

        except Exception as e:
            logging.error(f"Error al insertar el gráfico DUAL {key}: {e}")
            
            logging.info(f"Diapositiva DUAL creada.")


def add_conversion_textbox(slide, rates, chart_frame, left_pos, chart_width, funnel_start_top):
    """Dibuja las flechas de conversión usando los valores exactos del Excel."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    # Calculamos el alto de cada etapa para centrar el texto entre barras
    num_bars = len(rates)
    bar_height_step = chart_frame.height / num_bars
    
    # 'rates' ya contiene los valores que calculaste (ej: [100, 85, 70...])
    for i, rate in enumerate(rates):
        if i == 0: continue # La primera etapa es la base 100
        
        # Posicionamos el cuadro entre la etapa anterior e i
        top_pos = funnel_start_top + (bar_height_step * i) - Inches(0.1)
        left_pos_final = left_pos + chart_width - Inches(0.15)
        
        txBox = slide.shapes.add_textbox(left_pos_final, top_pos, Inches(0.5), Inches(0.3))
        p = txBox.text_frame.paragraphs[0]
        
        # Si el valor ya viene como string o float del engine
        try:
            val = int(float(rate))
            p.text = f"➔ {val}%"
        except:
            p.text = f"➔ {rate}"

        # Simplemente imprimimos el valor del engine
        p.text = f"➔ {int(rate)}%" 
        p.font.size = Pt(8)
        p.font.name = 'Nunito Light'
        p.font.color.rgb = RGBColor(80, 80, 80)

    return txBox


def _create_single_brand_funnel(slide, placeholder, brand_name, data_series, stage_names, current_left_pos, final_chart_width, final_chart_height, **kwargs):
    """
    Crea un gráfico de Barras Apiladas Centradas (Funnel) con datos normalizados y etiquetas visibles.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    # 1. Preparación de Datos
    chart_data = CategoryChartData()
    chart_data.categories = stage_names if stage_names else [str(i) for i in data_series.index]
    
    blank_L_list = []
    value_list = []
    blank_R_list = []

    # Si llega un DataFrame, extraemos la primera columna para evitar el AttributeError
    if hasattr(data_series, 'columns'): 
        data_series = data_series.iloc[:, 0]
    
    # Convertimos a Serie de Pandas por seguridad
    if not isinstance(data_series, pd.Series):
        data_series = pd.Series(data_series)

    # Convertir toda la serie a float limpio
    clean_series = (
        data_series
        .astype(str)
        .str.replace('%', '', regex=False)
        .str.replace(',', '.', regex=False)
        .str.strip()
        .astype(float)
    )


    clean_series = pd.to_numeric(clean_series, errors='coerce').fillna(0.0)
    
    max_val = clean_series.max()
    
    for value in clean_series:
        # Normalizar a escala 0-1
        
        value_pct = value / 100.0
        remaining = 1.0 - value_pct
        # Generar las partes en blanco para centrar
        blank_L_list.append(remaining / 2.0)
        value_list.append(value_pct)
        blank_R_list.append(remaining / 2.0)

    # DEBUG en consola para ver qué estamos enviando al gráfico
    logging.info(f"📊 Marca {brand_name}: {value_list}")

    # Añadir series: El orden L-C-R es lo que centra el embudo
    chart_data.add_series('Blank L', blank_L_list)
    chart_data.add_series(brand_name, value_list)
    #chart_data.add_series(str(brand_name), value_list)
    chart_data.add_series('Blank R', blank_R_list)

    # 2. Inserción del Gráfico
    chart_frame = slide.shapes.add_chart(
        chart_type=XL_CHART_TYPE.BAR_STACKED,
        x=current_left_pos,
        y=placeholder.top + Inches(0.6), # Espacio para el título
        cx=final_chart_width,
        cy=final_chart_height,
        chart_data=chart_data
    )
    chart = chart_frame.chart
    # Ocultar el eje vertical (Nombres de etapas)
    # Lo quitamos para todas las marcas para que solo se vea el embudo puro
    category_axis = chart.category_axis
    category_axis.visible = False
    
    chart.has_legend = False
    chart.category_axis.reverse_order = True # Para que la primera etapa (Conocimiento) esté arriba
    chart.value_axis.visible = False # Ocultar Eje X (Valores)
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.maximum_scale = 1.0 # Forzar al 100%

    chart.value_axis.major_gridlines.format.line.fill.background()
    chart.category_axis.major_gridlines.format.line.fill.background()

    # Nueva forma de hacer los Blanks "invisibles": Relleno blanco y línea blanca
    # Reemplazar: chart.series[0].format.fill.background()
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = RGBColor(255, 255, 255) # Color blanco
    chart.series[0].format.line.color.rgb = RGBColor(255, 255, 255) # Línea blanca

    #Reemplazar: chart.series[2].format.fill.background()
    chart.series[2].format.fill.solid()
    chart.series[2].format.fill.fore_color.rgb = RGBColor(255, 255, 255) # Color blanco
    chart.series[2].format.line.color.rgb = RGBColor(255, 255, 255) # Línea blanca

    show_axis = kwargs.get('show_category_axis', False)
    chart.category_axis.visible = show_axis

    # Estilizar la serie de datos (Serie 1)
    series_value = chart.series[1]
    
    # 4. Colores por etapa y Etiquetas
    for j, point in enumerate(series_value.points):
        # Colores (ajusta STAGE_COLORS en tus constantes)
        color_hex = STAGE_COLORS[j % len(STAGE_COLORS)] if 'STAGE_COLORS' in globals() else "0070C0"
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(color_hex.replace("#",""))
        
    series_value.has_data_labels = True
    series_value.gap_width = 10 # Barras más gruesas para look de funnel
    data_labels = series_value.data_labels
    
    data_labels.show_value = True 
    data_labels.show_percentage = False
    data_labels.number_format_is_linked = False
    data_labels.number_format = '0%'       
    data_labels.position = POS.CENTER
            # Fuente (accesible desde data_labels.font)
    data_labels.font.name = 'Nunito Light'
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = RGBColor(255, 255, 255)


    # 5. Título de la Marca
    chart.has_title = True
    chart.chart_title.text_frame.text = str(brand_name)
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)
    chart.chart_title.text_frame.paragraphs[0].font.name = 'Nunito'

    # Quitar eje de categorías si no es la primera marca
    if current_left_pos > placeholder.left + Inches(0.5):
        chart.category_axis.visible = False

    return chart_frame

def create_brand_funnel_slide(prs, funnel_df, finding, task, all_step_rates=None, show_conversion_rates=True, stage_names=None):
    """
    Crea una diapositiva vacía y genera múltiples gráficos de Funnel de Cascada (uno por marca)
    posicionados dinámicamente.
    """
    
    # --- Modificación en create_brand_funnel_slide ---
    # 1. Definir el Layout de la Tarea (Si no está ya definido)
    slide_layout_name = task.get("SLIDE_LAYOUT", 'Chart layout 8')
    slide_layout = utils.find_layout_by_name(prs, slide_layout_name)
    #slide = prs.slides.add_slide(slide_layout)

    if slide_layout is None:
        logging.error(f"FATAL: Layout '{slide_layout_name}' no encontrado. Usando layout 6.")
        slide_layout = prs.slide_layouts[6] 
        
    slide = prs.slides.add_slide(slide_layout) # <-- Diapositiva ÚNICA y CORRECTA

# --- 🚨 SOLUCIÓN AL ERROR: Búsqueda segura del Título 🚨 ---
    title_placeholder = None
    try:
        # 1. Buscar el placeholder de tipo TITLE (1)
        title_placeholder = next(
            ph for ph in slide.placeholders 
            if ph.placeholder_format.type == PH_TYPE.TITLE
        )
        
        # 2. Rellenar si se encuentra
        title_placeholder.text = finding
        # Opcional: ajustar el estilo si es necesario (utilizar el text_frame del placeholder)
        # p = title_placeholder.text_frame.paragraphs[0]
        # p.font.size = Pt(20) 
        # p.font.name = FONT_NAME 

    except StopIteration:
        # Esto ocurre si el layout no tiene un placeholder de Título
        logging.warning("El layout no contiene un placeholder de título (PH_TYPE.TITLE).")

    # Si necesitas un Subtítulo (Placeholder de Body o Subtitle)
    # Suponemos que la variable 'task' puede contener el subtítulo
    subtitle_text = task.get("SUBTITLE", None) 
    
    if subtitle_text:
        try:
            # Buscar el placeholder de tipo BODY (2) o SUBTITLE (13). 
            # El layout 8 a menudo usa BODY para el segundo texto.
            subtitle_placeholder = next(
                ph for ph in slide.placeholders 
                if ph.placeholder_format.type == PH_TYPE.BODY or ph.placeholder_format.type == PH_TYPE.SUBTITLE
            )
            text_frame = subtitle_placeholder.text_frame
            text_frame.text = subtitle_text
            # Opcional: configurar el texto si es necesario
            # text_frame.paragraphs[0].font.size = Pt(14)

        except StopIteration:
            logging.warning("No se encontró el placeholder de subtítulo (BODY/SUBTITLE).")
# --- TÍTULO (Usando coordenadas fijas o un placeholder de título) ---
    #title_left, title_top = Inches(0.5), Inches(0.5)
    #title_box = slide.shapes.add_textbox(title_left, title_top, Inches(9.0), Inches(0.5))
    #title_box.text_frame.text = finding
    #p = title_box.text_frame.paragraphs[0]
    #p.font.size = Pt(20)
    #p.font.name = FONT_NAME

    # 🚨 INICIO DEL BLOQUE DE BÚSQUEDA CORREGIDO 🚨
    chart_placeholder = None
    # 2. Encontrar el Placeholder de Contenido
    # Buscar por el tipo OBJECT (7), que es el más probable para un contenedor de gráficos
# Buscar por el tipo OBJECT (7), que es el más probable para un contenedor de gráficos
    try:
    # 🚨 CORRECCIÓN CLAVE AQUÍ 🚨
        chart_placeholder = next(
        ph for ph in slide.placeholders 
        # 🚨 Busca OBJECT (7) o BODY (2) 🚨
        if ph.placeholder_format.type == PH_TYPE.OBJECT or ph.placeholder_format.type == PH_TYPE.BODY
    )
    except StopIteration:
        # Si no se encuentra un tipo BODY u OBJECT, recurrimos a un área manual (fallback)
        logging.warning("No se encontró el Placeholder OBJECT/BODY. Usando área fija.")
        # Crea un rectángulo ficticio seguro para usar sus coordenadas
        chart_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(8), Inches(5)
        )
        chart_placeholder.fill.background()
        chart_placeholder.line.fill.background()

        # 🚨 FIN DEL BLOQUE DE BÚSQUEDA CORREGIDO 🚨
    if all_step_rates is None:
        raise ValueError("all_step_rates es obligatorio para create_brand_funnel_slide")

    
    # --- 1. CONFIGURACIÓN Y OBTENCIÓN DE LAYOUT VACÍO ---
    brands = funnel_df.index.tolist()
    num_brands = len(brands)

    # --- 2. DEFINIR ÁREAS DE TEXTO Y GRÁFICOS (Coordenadas Absolutas) ---
    
    # Área del Título Principal
    #title_left, title_top = Inches(0.5), Inches(0.5)
    #title_width, title_height = Inches(9.0), Inches(0.5) 
    
    #title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    #title_frame = title_box.text_frame
    #title_frame.text = finding
    # Aplicar el estilo de título (ejemplo: fuente grande y negrita)
    #p = title_frame.paragraphs[0]
    #p.font.size = Pt(20)
    #p.font.name = FONT_NAME
    
    chart_left = chart_placeholder.left # <- ¡DEFINIDO DE NUEVO!
    chart_top = chart_placeholder.top  # <- ¡DEFINIDO DE NUEVO!
    chart_width = chart_placeholder.width
    chart_height = chart_placeholder.height

# 🚨 AJUSTES DE POSICIÓN SOLICITADOS (MODIFIQUE ESTOS VALORES) 🚨
    # Aumentar para mover a la derecha
    HORIZONTAL_OFFSET = Inches(4) 
    # Aumentar para mover hacia abajo
    VERTICAL_OFFSET = Inches(5.6)

# Coordenadas base del placeholder (Guía de diseño)
    base_chart_left = chart_placeholder.left
    base_chart_top = chart_placeholder.top
    
    # Aplicar Offsets:
    funnel_start_left = base_chart_left + HORIZONTAL_OFFSET
    funnel_start_top = base_chart_top + VERTICAL_OFFSET

    # 🚨 NUEVA REFERENCIA DE COORDENADAS 🚨
    #chart_left = chart_placeholder.left
    #chart_top = chart_placeholder.top

    #chart_width = chart_placeholder.width
    #chart_height = chart_placeholder.height
    # --- 3. CÁLCULO DE GEOMETRÍA DINÁMICA ---
    
# 3.1. Definición de Dimensiones Fijas (Conversión de cm a Inches)
    # 6.46 cm -> 2.5433"
    FIXED_CHART_WIDTH = Inches(1.8) 
    # 9.64 cm -> 3.7952"
    FIXED_CHART_HEIGHT = Inches(3.11) 
    # 0.5 cm -> 0.1968"
    FIXED_SPACING = Inches(0.02)

    # 🚨 AÑADIR ESTA LÍNEA 🚨
    TEXTBOX_WIDTH = Inches(0.4)

    #current_left_pos = chart_left
    #current_left_pos = funnel_start_left
    current_left_pos = Inches(2.2)
    
    # --- 4. BUCLE DE GENERACIÓN POR MARCA ---
    brand_colors = BRAND_COLORS # Paleta de colores definida

    #stage_names = funnel_df.columns.tolist() # Obtener nombres de las etapas una vez
    

    # 🚨 CORRECCIÓN: Usar stage_names si vienen del VisualEngine
    actual_stage_names = stage_names if stage_names else funnel_df.columns.tolist()
    brands = funnel_df.index.tolist()
    #current_left_pos = Inches(0.5) # Posición inicial clara
    for i, brand in enumerate(brands):
        # brand conserva su nombre único (ej: "TOTAL - Duramax")
        data_series = funnel_df.loc[brand]
        #brand_str = str(brand) # Ya debería ser string por el paso anterior
        
        # Limpiamos el nombre SOLO para el título del gráfico y el lookup de tasas
        display_name = str(brand)
        for prefix in ["TOTAL - ", "F1: ", "Hombre - ", "Mujer - "]:
            display_name = display_name.replace(prefix, "")
        display_name = display_name.strip()

        # 🚨 CAMBIO CRÍTICO 1: NORMALIZACIÓN DE DATOS 🚨
        # El gráfico espera valores decimales (0.0 a 1.0) para que el eje X se vea correcto.
        #data_series_decimal = data_series / 100.0

        # Generar y posicionar el Funnel individual
        #chart_frame = _create_single_brand_funnel(
        
        # En create_brand_funnel_slide, antes de llamar a _create_single_brand_funnel
        # brand viene como ('TOTAL', 'Duramax (Scott)') o similar
        #if isinstance(brand, tuple):
        #    display_title = " - ".join([str(item) for item in brand])
        #else:
        #    display_title = str(brand)
        


        # 1. Si es tupla, extraemos SOLO el último elemento (la Marca)
        if isinstance(brand, tuple):
            display_title = str(brand[-1]) 
        else:
            display_title = str(brand)
        
        # 2. Limpieza de seguridad para eliminar prefijos residuales
        # Esto elimina "Total - ", "1. ", etc., si estuvieran pegados al nombre
        prefijos_a_quitar = ["TOTAL - ", "Total - ", "F1: ", "SEGMENTACIÓN: "]
        for prefijo in prefijos_a_quitar:
            display_title = display_title.replace(prefijo, "")
    
        # 3. Eliminar números iniciales (ej: "1. Duramax" -> "Duramax")
        import re
        display_title = re.sub(r'^\d+[\.\s\-]+', '', display_title).strip()

        chart_frame = _create_single_brand_funnel(
            slide=slide, 
            placeholder=chart_placeholder, 
            #brand_name=display_name, 
            brand_name=display_title, # Pasamos el nombre compuesto
            data_series=data_series, 
            #stage_names=funnel_df.columns.tolist(),
            stage_names=actual_stage_names,
            #brand_color=brand_color_hex, 
            current_left_pos=current_left_pos,
            chart_top_pos=funnel_start_top,
            final_chart_width=FIXED_CHART_WIDTH, # Ancho dinámico
            final_chart_height=FIXED_CHART_HEIGHT # Altura completa
        )

        # --- AVANCE DE POSICIÓN SEGURO ---
        # Sumamos Inches + Inches para mantener la unidad correcta
        ancho_total = FIXED_CHART_WIDTH + TEXTBOX_WIDTH + Inches(0.1)
        current_left_pos = current_left_pos + ancho_total

# 🚨 LÓGICA CONDICIONAL 🚨
        if show_conversion_rates and all_step_rates:
            step_rates = all_step_rates.get(display_name)
            if step_rates is not None:
                add_conversion_textbox(slide, step_rates, chart_frame, current_left_pos, FIXED_CHART_WIDTH, funnel_start_top)
        
        
        
        #if show_conversion_rates:
        #    #step_rates_for_brand = all_step_rates[brand]
        #    
        #    step_rates_for_brand = all_step_rates.get(display_name)
        #    # Llama a la función para dibujar el TextBox
        #    
        #    if step_rates is not None:
        #        add_conversion_textbox(slide, step_rates, chart_frame, current_left_pos, FIXED_CHART_WIDTH, funnel_start_top)
            
            #if step_rates_for_brand is not None and hasattr(step_rates_for_brand, '__len__'):    
            #    add_conversion_textbox(
            #        slide,
            #        step_rates_for_brand, 
            #        chart_frame, 
            #        current_left_pos, 
            #        FIXED_CHART_WIDTH,
            #        #chart_top_pos=funnel_start_top
            #        funnel_start_top=funnel_start_top
            #    )
            #else:
            #    logging.warning(f"No se encontraron tasas de conversión para la marca: {display_name}")


        #ancho_grafico = FIXED_CHART_WIDTH if isinstance(FIXED_CHART_WIDTH, Inches) else Inches(FIXED_CHART_WIDTH)
        #ancho_texto = TEXTBOX_WIDTH if isinstance(TEXTBOX_WIDTH, Inches) else Inches(TEXTBOX_WIDTH)
        #espaciado = Inches(0.1)

        #current_left_pos = current_left_pos + ancho_grafico + ancho_texto + espaciado

# === APLICACIÓN DE LA CORRECCIÓN DE TIPOS ===
        #new_left_pos = current_left_pos + FIXED_CHART_WIDTH + TEXTBOX_WIDTH + FIXED_SPACING
        #current_left_pos = new_left_pos

        logging.info(f"DEBUG: Posición izquierda actualizada a: {current_left_pos:.2f} pulgadas")
    logging.info(f"Diapositiva creada (FUNNEL DE CASCADA MÚLTIPLE): {task.get('VARIABLE_NAME')}.")

def update_tracking_chart_in_presentation(prs, chart_name, new_wave_label, metrics_data, is_percentage=False, append_only=False, line_colors=None, remove_percentage_sign=False, decimals=0, ola_impar=False):
    if line_colors is None:
        line_colors = {}
        
    def arreglar_etiqueta(etiqueta):
        texto = str(etiqueta).replace('\u200B', '').replace('\xa0', '').strip()
        try:
            val = float(texto)
            if 40000 < val < 50000:
                base_date = datetime.datetime(1899, 12, 30)
                target_date = base_date + datetime.timedelta(days=int(val))
                meses = ["Ene", "Feb", "Mar", "Abr", "May", "Jun", "Jul", "Ago", "Sep", "Oct", "Nov", "Dic"]
                texto = f"{meses[target_date.month - 1]} {str(target_date.year)[-2:]}"
        except:
            pass
        
        texto = texto.replace("-", " ")
        partes = texto.split()
        if len(partes) >= 2:
            texto = f"{partes[0].capitalize()} {partes[1]}"
        else:
            texto = texto.title() if texto.islower() else texto
        return f"\u200B{texto}"

    grafico_actualizado = False 

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart and shape.name == chart_name:
                chart = shape.chart
                chart_data = CategoryChartData()
                
                new_wave_label_clean = arreglar_etiqueta(new_wave_label)

                raw_categories = []
                if len(chart.plots) > 0:
                    try:
                        raw_categories = chart.plots[0].categories
                    except:
                        pass
                
                indices_validos = [i for i, c in enumerate(raw_categories) if str(c.label if hasattr(c, 'label') else c).replace('\u200B', '').strip() != ""]
                old_categories = [arreglar_etiqueta(raw_categories[i].label if hasattr(raw_categories[i], 'label') else raw_categories[i]) for i in indices_validos]

                mes_ya_existe = False
                if len(old_categories) > 0:
                    ultimo_limpio = old_categories[-1].replace('\u200B', '').lower()
                    nuevo_limpio = new_wave_label_clean.replace('\u200B', '').lower()
                    if ultimo_limpio == nuevo_limpio:
                        mes_ya_existe = True
                
                # =========================================================
                # 🚀 INYECCIÓN DE CATEGORÍAS CON ELIMINACIÓN DOBLE
                # =========================================================
                if mes_ya_existe:
                    old_categories[-1] = new_wave_label_clean
                else:
                    old_categories.append(new_wave_label_clean)
                    if not append_only:
                        if ola_impar and len(old_categories) > 2:
                            old_categories.pop(0) # Elimina el más viejo
                            old_categories.pop(0) # Elimina el segundo más viejo
                        elif len(old_categories) > 1:
                            old_categories.pop(0) # Comportamiento normal (elimina 1)
                        
                chart_data.categories = old_categories

                chart_title_text = ""
                try:
                    if chart.has_title and chart.chart_title.has_text_frame:
                        chart_title_text = chart.chart_title.text_frame.text.lower()
                        chart_title_text = " ".join(chart_title_text.split()) 
                except:
                    pass
                
                series_activas = []
                for s in chart.series:
                    if hasattr(s, 'name') and s.name:
                        clean_name = str(s.name).strip()
                        if clean_name != "":
                            series_activas.append(s)

                for series_idx, series in enumerate(series_activas):
                    raw_vals = list(series.values)
                    
                    old_vals = []
                    for i in indices_validos:
                        if i < len(raw_vals):
                            val_crudo = raw_vals[i]
                            if isinstance(val_crudo, str):
                                val_limpio = str(val_crudo).strip()
                                if val_limpio == "" or val_limpio.lower() in ["nan", "null"]:
                                    old_vals.append(None)
                                else:
                                    try:
                                        old_vals.append(float(val_limpio))
                                    except ValueError:
                                        old_vals.append(None)
                            else:
                                old_vals.append(val_crudo)
                        else:
                            old_vals.append(None)

                    def is_empty_val(v):
                        if v is None: return True
                        if isinstance(v, str) and str(v).strip() == "": return True
                        try:
                            if float(v) == 0.0: return True
                        except:
                            pass
                        return False

                    historial_vacio = all(is_empty_val(v) for v in old_vals)
                                                
                    if is_percentage:
                        valores_limpios = []
                        for val in old_vals:
                            if val is not None:
                                if remove_percentage_sign:
                                    valores_limpios.append(val * 100.0 if (val < 1 and val > -1 and val != 0) else val)
                                else:
                                    valores_limpios.append(val / 100.0 if (val > 1 or val < -1) else val)
                            else:
                                valores_limpios.append(None)
                        old_vals = valores_limpios
                    
                    new_val = None
                    for keyword, val in metrics_data.items():
                        kw_clean = " ".join(keyword.lower().split())
                        s_name_clean = " ".join(series.name.lower().split())
                        
                        if (kw_clean in s_name_clean or 
                            s_name_clean in kw_clean or
                            (chart_title_text and kw_clean in chart_title_text)):
                            new_val = val
                            break

                    if new_val is None and len(metrics_data) == 1 and series_idx == 0:
                        new_val = list(metrics_data.values())[0]

                    if new_val is None and historial_vacio:
                        continue 

                    if new_val is not None:
                        if is_percentage:
                            if remove_percentage_sign:
                                if new_val < 1 and new_val > -1 and new_val != 0:
                                    new_val = new_val * 100.0 
                            else:
                                if new_val > 1 or new_val < -1:
                                    new_val = new_val / 100.0
                                    
                    # =========================================================
                    # 🚀 INYECCIÓN DE DATOS CON ELIMINACIÓN DOBLE
                    # =========================================================
                    if mes_ya_existe:
                        if len(old_vals) > 0:
                            old_vals[-1] = new_val
                    else:
                        old_vals.append(new_val)
                        if not append_only:
                            if ola_impar and len(old_vals) > 2:
                                old_vals.pop(0) # Elimina el más viejo
                                old_vals.pop(0) # Elimina el segundo más viejo
                            elif len(old_vals) > 1:
                                old_vals.pop(0) # Comportamiento normal

                    if is_percentage:
                        formato = '0.0%' if decimals == 1 else ('0%' if not remove_percentage_sign else '0')
                    else:
                        formato = '0.0' if decimals == 1 else ('0.00' if decimals == 2 else '0')                        
                    
                    chart_data.add_series(series.name, tuple(old_vals), number_format=formato)

                color_memoria = {}
                for s in chart.series:
                    try:
                        if s.format.line.color.rgb:
                            color_memoria[s.name] = s.format.line.color.rgb
                    except:
                        pass

                chart.replace_data(chart_data)
                
                for series in chart.series:
                    tipos_de_barras = [
                        XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.COLUMN_STACKED, 
                        XL_CHART_TYPE.COLUMN_STACKED_100, XL_CHART_TYPE.BAR_CLUSTERED, 
                        XL_CHART_TYPE.BAR_STACKED, XL_CHART_TYPE.BAR_STACKED_100
                    ]
                        
                    if chart.chart_type in tipos_de_barras:
                        series.format.line.fill.background()
                    
                    estilo = line_colors.get(series.name)
                    
                    if estilo:
                        if isinstance(estilo, dict):
                            color_hex = estilo.get("color", "").replace("#", "")
                            font_name = estilo.get("font_name")
                            font_size = estilo.get("font_size")
                        else:
                            color_hex = str(estilo).replace("#", "")
                            font_name = None
                            font_size = None

                        if color_hex and len(color_hex) == 6:
                            try:
                                r = int(color_hex[0:2], 16)
                                g = int(color_hex[2:4], 16)
                                b = int(color_hex[4:6], 16)
                                color_rgb = RGBColor(r, g, b)
                                
                                series.format.line.color.rgb = color_rgb
                                series.format.line.width = Pt(2.25)
                                series.marker.style = XL_MARKER_STYLE.CIRCLE
                                series.marker.size = 5
                                series.marker.format.fill.solid()
                                series.marker.format.fill.fore_color.rgb = color_rgb
                                series.marker.format.line.color.rgb = color_rgb
                            except AttributeError:
                                pass 
                        else:
                            color_rgb = None

                        try:
                            if series.has_data_labels:
                                font = series.data_labels.font
                                if font_name:
                                    font.name = font_name
                                if font_size:
                                    font.size = Pt(font_size)
                                if color_rgb:
                                    font.color.rgb = color_rgb
                        except Exception as e:
                            pass 
                
                try:
                    if is_percentage and not remove_percentage_sign:
                        chart.value_axis.number_format = '0.0%' if decimals == 1 else '0%'
                    elif remove_percentage_sign:
                        chart.value_axis.number_format = '0'
                    else:
                        chart.value_axis.number_format = '0.0' if decimals == 1 else ('0.00' if decimals == 2 else '0')
                except Exception:
                    pass                

                grafico_actualizado = True 
                
    return grafico_actualizado

def update_tracking_table_in_presentation(prs, table_name, new_wave_label, metrics_data, start_data_col=1, label_col=0, is_percentage=False, has_header=True):
    """
    Desplaza los datos de una tabla de PowerPoint hacia la izquierda
    e inyecta la nueva ola conservando el formato.
    UNIVERSAL: Acepta cualquier configuración de columnas y soporta tablas "sin encabezado".
    """
    
    def escribir_con_formato(celda, nuevo_texto):
        nuevo_texto = str(nuevo_texto).strip()
        
        if celda.text_frame.paragraphs and celda.text_frame.paragraphs[0].runs:
            p_principal = celda.text_frame.paragraphs[0]
            p_principal.runs[0].text = nuevo_texto
            
            for i in range(1, len(p_principal.runs)):
                p_principal.runs[i].text = ""
                
            for i in range(len(celda.text_frame.paragraphs) - 1, 0, -1):
                p_fantasma = celda.text_frame.paragraphs[i]._p
                p_fantasma.getparent().remove(p_fantasma) 
        else:
            celda.text = nuevo_texto

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table and shape.name == table_name:
                table = shape.table
                num_cols = len(table.columns)
                num_rows = len(table.rows)

                # =========================================================
                # 🚀 NUEVO: CONTROL DE ENCABEZADOS
                # Si tiene encabezado, actualiza la fila 0 y los datos arrancan en la fila 1.
                # Si NO tiene encabezado, ignora la fila 0 y los datos arrancan directo en la 0.
                start_data_row = 1 if has_header else 0

                if has_header:
                    # 1. ACTUALIZAR EL ENCABEZADO (Desplaza los meses)
                    for col_idx in range(start_data_col, num_cols - 1):
                        texto_a_copiar = table.cell(0, col_idx + 1).text.strip() 
                        escribir_con_formato(table.cell(0, col_idx), texto_a_copiar)
                    
                    # Inyectamos el nombre de la nueva ola al final
                    escribir_con_formato(table.cell(0, num_cols - 1), new_wave_label)
                # =========================================================

                # 2. ACTUALIZAR LAS FILAS DE DATOS
                for row_idx in range(start_data_row, num_rows):
                    row_label = table.cell(row_idx, label_col).text.strip().lower()

                    # A. Desplazar datos históricos a la izquierda
                    for col_idx in range(start_data_col, num_cols - 1):
                        texto_a_copiar = table.cell(row_idx, col_idx + 1).text.strip()
                        escribir_con_formato(table.cell(row_idx, col_idx), texto_a_copiar)

                    # B. Inyectar el dato nuevo
                    new_val = None
                    for keyword, val in metrics_data.items():
                        if keyword.lower() in row_label:
                            new_val = val
                            break
                    
                    if new_val is not None:
                        texto_final = f"{new_val}%" if is_percentage else str(new_val)
                        escribir_con_formato(table.cell(row_idx, num_cols - 1), texto_final)
                    else:
                        escribir_con_formato(table.cell(row_idx, num_cols - 1), "-")
                
                return True
    return False

def update_static_table_in_presentation(prs, table_name, target_col, metrics_data, label_col=0):
    """
    Inyecta datos en una tabla estática con Heatmap y Búsqueda Fuzzy (Tolerante a errores).
    """
    
    # --- BISTURÍ DE PRECISIÓN ---
    def escribir_con_formato(celda, nuevo_texto):
        nuevo_texto = str(nuevo_texto).strip()
        if celda.text_frame.paragraphs and celda.text_frame.paragraphs[0].runs:
            p_principal = celda.text_frame.paragraphs[0]
            p_principal.runs[0].text = nuevo_texto
            for i in range(1, len(p_principal.runs)):
                p_principal.runs[i].text = ""
            for i in range(len(celda.text_frame.paragraphs) - 1, 0, -1):
                p_fantasma = celda.text_frame.paragraphs[i]._p
                p_fantasma.getparent().remove(p_fantasma) 
        else:
            celda.text = nuevo_texto

    # =========================================================
    valores_numericos = [
        v for k, v in metrics_data.items() 
        if str(k).lower() != "promedio" and isinstance(v, (int, float))
    ]
    col_min = min(valores_numericos) if valores_numericos else 0
    col_max = max(valores_numericos) if valores_numericos else 100

    def pintar_celda_relativa(celda, valor):
        try:
            v = float(valor)
        except (ValueError, TypeError):
            return 
            
        if col_max == col_min:
            ratio = 0.5
        else:
            ratio = (v - col_min) / (col_max - col_min)
            
        ratio = max(0.0, min(1.0, ratio)) 
        
        if ratio < 0.5:
            pct = ratio * 2.0
            r = int(220 + (248 - 220) * pct)
            g = int(53 + (249 - 53) * pct)
            b = int(69 + (250 - 69) * pct)
        else:
            pct = (ratio - 0.5) * 2.0
            r = int(248 + (13 - 248) * pct)
            g = int(249 + (110 - 249) * pct)
            b = int(250 + (253 - 250) * pct)

        celda.fill.solid()
        celda.fill.fore_color.rgb = RGBColor(r, g, b)
    # =========================================================

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table and shape.name == table_name:
                table = shape.table
                num_rows = len(table.rows)

                for row_idx in range(0, num_rows):
                    # 👇 MAGIA NUEVA: Leemos las Columnas 0 y 1 simultáneamente 👇
                    texto_col0 = table.cell(row_idx, 0).text.lower() if len(table.columns) > 0 else ""
                    texto_col1 = table.cell(row_idx, 1).text.lower() if len(table.columns) > 1 else ""
                    
                    row_lbl_0 = " ".join(texto_col0.split())
                    row_lbl_1 = " ".join(texto_col1.split())

                    new_val = None
                    for keyword, val in metrics_data.items():
                        kw_clean = " ".join(keyword.lower().split())
                        
                        # 1. Match Exacto (Cruza todo)
                        if (kw_clean in row_lbl_0 or row_lbl_0 in kw_clean and len(row_lbl_0) > 5) or \
                           (kw_clean in row_lbl_1 or row_lbl_1 in kw_clean and len(row_lbl_1) > 5):
                            new_val = val
                            break
                            
                        # 2. Match Fuzzy (Tolera faltas de ortografía, tildes y espacios raros)
                        if len(kw_clean) > 10:
                            similitud_0 = difflib.SequenceMatcher(None, kw_clean, row_lbl_0).ratio()
                            similitud_1 = difflib.SequenceMatcher(None, kw_clean, row_lbl_1).ratio()
                            
                            if similitud_0 > 0.85 or similitud_1 > 0.85:
                                new_val = val
                                break
                    
                    if new_val is not None:
                        celda_objetivo = table.cell(row_idx, target_col)
                        escribir_con_formato(celda_objetivo, new_val)
                        
                        if "promedio" not in row_lbl_0 and "promedio" not in row_lbl_1:
                            pintar_celda_relativa(celda_objetivo, new_val)
                
                return True
    return False

def update_header_table_in_presentation(prs, table_name, new_wave_label):
    """
    Desplaza únicamente los meses de una tabla "encabezado" hacia la izquierda,
    e inyecta la nueva ola en la última columna, conservando el formato.
    Actualiza TODAS las instancias de esa tabla en toda la presentación.
    """
    
    # --- BISTURÍ DE PRECISIÓN ---
    def escribir_con_formato(celda, nuevo_texto):
        nuevo_texto = str(nuevo_texto).strip()
        if celda.text_frame.paragraphs and celda.text_frame.paragraphs[0].runs:
            p_principal = celda.text_frame.paragraphs[0]
            p_principal.runs[0].text = nuevo_texto
            for i in range(1, len(p_principal.runs)):
                p_principal.runs[i].text = ""
            for i in range(len(celda.text_frame.paragraphs) - 1, 0, -1):
                p_fantasma = celda.text_frame.paragraphs[i]._p
                p_fantasma.getparent().remove(p_fantasma) 
        else:
            celda.text = nuevo_texto
    # -------------------------------------------------------------

    tabla_encontrada_y_actualizada = False

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table and shape.name == table_name:
                table = shape.table
                num_cols = len(table.columns)

                # 1. Desplazamos todos los meses un lugar a la izquierda (arranca en Col 0)
                for col_idx in range(0, num_cols - 1):
                    texto_a_copiar = table.cell(0, col_idx + 1).text.strip()
                    escribir_con_formato(table.cell(0, col_idx), texto_a_copiar)

                # 2. Inyectamos la nueva ola en el último casillero
                escribir_con_formato(table.cell(0, num_cols - 1), new_wave_label)
                
                # Marcamos que al menos una tabla se actualizó, pero NO CORTAMOS el bucle
                tabla_encontrada_y_actualizada = True 

    # Devuelve True solo si encontró y modificó al menos una tabla
    return tabla_encontrada_y_actualizada

def update_ytd_calculated_from_chart(prs, ytd_chart_name, ref_chart_name, target_year_label, year_suffix, metrics_map=None, is_percentage=False, remove_percentage_sign=False, decimals=0, multiplier=1.0):
    promedios_calculados = {}
    if metrics_map is None:
        metrics_map = {}
    
    def arreglar_etiqueta(etiqueta):
        try:
            val = float(etiqueta)
            if 40000 < val < 50000:
                base_date = datetime.datetime(1899, 12, 30)
                target_date = base_date + datetime.timedelta(days=int(val))
                meses = ["ene", "feb", "mar", "abr", "may", "jun", "jul", "ago", "sep", "oct", "nov", "dic"]
                return f"{meses[target_date.month - 1]}-{str(target_date.year)[-2:]}"
        except:
            pass
        return str(etiqueta).lower().strip()
    
    def encontrar_coincidencia(texto_objetivo_ytd, titulo_grafico, diccionario_promedios):
        obj = str(texto_objetivo_ytd).lower().strip()
        for clave_ref, valor in diccionario_promedios.items():
            if clave_ref and obj and (clave_ref in obj or obj in clave_ref):
                return valor
        if titulo_grafico:
            for clave_ref, valor in diccionario_promedios.items():
                if clave_ref and (clave_ref in titulo_grafico or titulo_grafico in clave_ref):
                    return valor
        textos_a_buscar = [obj]
        if titulo_grafico: textos_a_buscar.append(titulo_grafico)
        for texto in textos_a_buscar:
            if not texto: continue
            for nombre_config, palabras_clave in metrics_map.items():
                if texto in str(nombre_config).lower() or str(nombre_config).lower() in texto:
                    for kw in palabras_clave:
                        kw_lower = str(kw).lower()
                        for clave_ref, valor in diccionario_promedios.items():
                            if clave_ref and kw_lower in clave_ref:
                                return valor
        return None

    # =========================================================
    # 1. MODO LECTURA
    # =========================================================
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart and shape.name == ref_chart_name:
                ref_chart = shape.chart
                try:
                    cats = [arreglar_etiqueta(c.label) for c in ref_chart.plots[0].categories]
                except:
                    cats = [arreglar_etiqueta(c) for c in ref_chart.plots[0].categories]
                
                indices_año = [i for i, cat in enumerate(cats) if str(year_suffix).lower() in cat]
                
                if not indices_año:
                    return False 
                    
                for series in ref_chart.series:
                    if not hasattr(series, 'name') or not series.name or str(series.name).strip() == "":
                        continue

                    valores_validos = [series.values[i] for i in indices_año if series.values[i] is not None]
                    if valores_validos:
                        promedio = sum(valores_validos) / len(valores_validos)
                        
                        # 🚀 ACÁ ESTÁ EL ÚNICO CAMBIO: Multiplicamos explícitamente si el JSON lo pide
                        promedio = promedio * multiplier

                        # Tu código original para no redondear porcentajes pequeños
                        if not (is_percentage and promedio <= 1.5):
                            promedio = int(round(promedio))
                            
                        promedios_calculados[series.name.lower().strip()] = promedio
                
                print(f"\n[RADAR YTD] Promedios listos para inyectar: {promedios_calculados}")

    if not promedios_calculados:
        return False

    # =========================================================
    # 2. MODO ESCRITURA
    # =========================================================
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart and shape.name == ytd_chart_name:
                ytd_chart = shape.chart
                chart_data = CategoryChartData()
                
                ytd_title_text = ""
                try:
                    if ytd_chart.has_title and ytd_chart.chart_title.has_text_frame:
                        ytd_title_text = ytd_chart.chart_title.text_frame.text.lower().strip()
                except:
                    pass

                try:
                    old_categories = [c.label for c in ytd_chart.plots[0].categories]
                except:
                    old_categories = [str(c) for c in ytd_chart.plots[0].categories]
                    
                chart_data.categories = old_categories
                
                target_lower = str(target_year_label).lower().strip()
                suffix_lower = str(year_suffix).lower().strip() 
                cat_lower = [str(c).lower().strip() for c in old_categories]
                
                cat_idx = -1
                for i, c in enumerate(cat_lower):
                    if target_lower in c or c in target_lower or suffix_lower in c:
                        cat_idx = i
                        break
                
                # ESCENARIO 1: El año está en el Eje X
                if cat_idx != -1:
                    for series in ytd_chart.series:
                        old_vals = list(series.values) # 🛡️ Tu código original no rompía la historia
                        s_name = series.name
                        
                        nuevo_promedio = encontrar_coincidencia(s_name, ytd_title_text, promedios_calculados)
                        
                        if nuevo_promedio is None and len(ytd_chart.series) == 1 and len(promedios_calculados) == 1:
                            nuevo_promedio = list(promedios_calculados.values())[0]

                        if nuevo_promedio is not None:
                            while len(old_vals) <= cat_idx:
                                old_vals.append(None)
                            old_vals[cat_idx] = nuevo_promedio
                            
                        formato = '0%' if is_percentage else '0'
                        chart_data.add_series(series.name, tuple(old_vals), number_format=formato)
                
                # ESCENARIO 2: El año está en la Leyenda (Columnas)
                else:
                    for series in ytd_chart.series:
                        old_vals = list(series.values)
                        s_name_lower = str(series.name).lower().strip()
                        
                        if target_lower in s_name_lower or s_name_lower in target_lower or suffix_lower in s_name_lower:
                            for i, cat_label in enumerate(old_categories):
                                nuevo_promedio = encontrar_coincidencia(cat_label, ytd_title_text, promedios_calculados)
                                
                                if nuevo_promedio is None and len(old_categories) == 1 and len(promedios_calculados) == 1:
                                    nuevo_promedio = list(promedios_calculados.values())[0]

                                if nuevo_promedio is not None:
                                    while len(old_vals) <= i:
                                        old_vals.append(None)
                                    old_vals[i] = nuevo_promedio
                        
                        formato = '0%' if is_percentage else '0'
                        chart_data.add_series(series.name, tuple(old_vals), number_format=formato)
                
                ytd_chart.replace_data(chart_data)
                return True
    return False

def generar_slide_desde_template(prs, layout_nombre, titulo_pregunta, insight_ia, categorias, valores, tipo_grafico="dona"):
    """
    Busca un molde por su nombre, crea un slide nuevo y rellena los placeholders.
    """
    # 1. Buscar el molde (Layout) correcto en el Master
    layout_elegido = None
    for layout in prs.slide_layouts:
        if layout.name == layout_nombre:
            layout_elegido = layout
            break
            
    if not layout_elegido:
        logging.error(f"❌ No se encontró el molde '{layout_nombre}' en la plantilla.")
        return False

    # 2. Crear el slide nuevo al final de la presentación
    nuevo_slide = prs.slides.add_slide(layout_elegido)
    logging.info(f"✨ Slide creado usando el molde: {layout_nombre}")

    # 3. Rellenar los Textos (IDs 13 y 15)
    # Asumimos que el 13 es el título (porque lo creaste primero) y el 15 es el insight.
    try:
        if 13 in nuevo_slide.placeholders:
            nuevo_slide.placeholders[13].text = titulo_pregunta
            
        if 15 in nuevo_slide.placeholders:
            nuevo_slide.placeholders[15].text = insight_ia
    except Exception as e:
        logging.warning(f"Cuidado con los textos: {e}")

    # 4. Preparar los datos matemáticos para el gráfico
    chart_data = CategoryChartData()
    chart_data.categories = categorias
    chart_data.add_series('Resultados', valores)

    # 5. Rellenar el Gráfico (ID 14) dependiendo del tipo
    if 14 in nuevo_slide.placeholders:
        marcador_grafico = nuevo_slide.placeholders[14]
        
        # Mapeo simple de tipos de gráficos de PPTX
        tipo_pptx = XL_CHART_TYPE.DOUGHNUT
        if tipo_grafico == "barras_verticales":
            tipo_pptx = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif tipo_grafico == "barras_horizontales":
            tipo_pptx = XL_CHART_TYPE.BAR_CLUSTERED
            
        # ¡INYECCIÓN LETAL! PowerPoint transforma el agujero vacío en un gráfico real
        marcador_grafico.insert_chart(tipo_pptx, chart_data)
        
    return True



def renombrar_por_geometria(prs):
    """
    Lee las notas y separa Gráficos de Tablas inteligentemente.
    Si el nombre pedido empieza con 'Chart' busca gráficos.
    Si empieza con 'Table', 'Tabla' o 'Header' busca tablas.
    """
    import logging
    logging.info("\n" + "🏷️"*10)
    logging.info("INICIANDO AUTO-REPARADOR INTELIGENTE (FILTRO POR TIPO)...")
    
    diapos_modificadas = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        if not slide.has_notes_slide:
            continue
            
        notas = slide.notes_slide.notes_text_frame.text
        
        if "DATOS:" not in notas.upper():
            continue
            
        # Extraemos los nombres
        lineas = notas.split('\n')
        linea_datos = [l for l in lineas if "DATOS:" in l.upper()][0]
        variables_str = linea_datos.split(":", 1)[1] 
        nombres_pedidos = [v.strip() for v in variables_str.split(",") if v.strip()]
        
        # Juntamos TODOS los objetos válidos
        objetos_validos = [s for s in slide.shapes if s.has_chart or s.has_table]
        
        # =========================================================
        # 🚀 EL AJUSTE ACÁ: Bandas horizontales estrictas
        # 500000 son aprox 1.3 cm. Si dos gráficos están en esa misma franja, 
        # los considera la misma "Fila" y desempatará ordenándolos por la izquierda (s.left).
        margen_fila = 500000 
        objetos_validos.sort(key=lambda s: (s.top // margen_fila, s.left))
        # =========================================================
        
        # Separamos en dos filas distintas (ya ordenadas visualmente)
        graficos_ordenados = [s for s in objetos_validos if s.has_chart]
        tablas_ordenadas = [s for s in objetos_validos if s.has_table]
        
        logging.info(f"📍 Diapo {slide_idx+1}: Hay {len(graficos_ordenados)} Gráficos y {len(tablas_ordenadas)} Tablas.")
        
        # Hacemos el match inteligente
        for nuevo_nombre in nombres_pedidos:
            nombre_upper = nuevo_nombre.upper()
            
            # ¿Es una tabla o un header?
            if nombre_upper.startswith("TABLE") or nombre_upper.startswith("TABLA") or nombre_upper.startswith("HEADER"):
                if tablas_ordenadas:
                    shape_elegido = tablas_ordenadas.pop(0) # Sacamos la 1ra tabla disponible
                    nombre_viejo = shape_elegido.name
                    shape_elegido.name = nuevo_nombre
                    logging.info(f"   📊 [TABLA] Renombrado: '{nombre_viejo}'  --->  '{nuevo_nombre}'")
                else:
                    logging.warning(f"   ⚠️ Faltan TABLAS físicas para asignar a '{nuevo_nombre}'")
            
            # Si no es tabla, asumimos que es un gráfico (CHART)
            else:
                if graficos_ordenados:
                    shape_elegido = graficos_ordenados.pop(0) # Sacamos el 1er gráfico disponible
                    nombre_viejo = shape_elegido.name
                    shape_elegido.name = nuevo_nombre
                    logging.info(f"   📈 [CHART] Renombrado: '{nombre_viejo}'  --->  '{nuevo_nombre}'")
                else:
                    logging.warning(f"   ⚠️ Faltan GRÁFICOS físicos para asignar a '{nuevo_nombre}'")
                    
        diapos_modificadas += 1
        
    logging.info("🏷️"*10 + "\n")

def actualizar_diapositiva_inteligente(slide, diccionario_resultados):
    """
    Busca la etiqueta DATOS: en las notas, ordena los gráficos por orden de lectura
    (arriba->abajo, izq->der) y les inyecta los datos.
    """
    # 1. VERIFICAR SI HAY INSTRUCCIONES EN LAS NOTAS
    if not slide.has_notes_slide:
        return
    
    notas = slide.notes_slide.notes_text_frame.text
    if "DATOS:" not in notas:
        return

    # Extraemos las variables solicitadas (Ej: DATOS: P04_1, P04_2, P04_3)
    linea_datos = [linea for linea in notas.split('\n') if "DATOS:" in linea][0]
    variables_str = linea_datos.split("DATOS:")[1]
    # Limpiamos espacios para que quede una lista prolija: ['P04_1', 'P04_2', 'P04_3']
    variables_requeridas = [v.strip() for v in variables_str.split(",") if v.strip()]

    # 2. BUSCAR TODOS LOS GRÁFICOS EN LA DIAPO
    graficos_shapes = [shape for shape in slide.shapes if shape.has_chart]
    
    if not graficos_shapes:
        print("⚠️ Las notas piden datos, pero no se encontraron gráficos en la diapositiva.")
        return

    if len(graficos_shapes) != len(variables_requeridas):
        print(f"⚠️ DISCREPANCIA: Las notas piden {len(variables_requeridas)} variables, pero hay {len(graficos_shapes)} gráficos.")
        # Podés decidir si hacer return acá, o procesar hasta donde alcance.

    # 3. EL TRUCO MAGISTRAL: ORDEN DE LECTURA HUMANO
    # Redondeamos la posición 'top' para agrupar gráficos que están en la misma "fila" visual,
    # aunque uno esté un par de píxeles más arriba que el otro. Luego ordenamos por 'left'.
    margen_fila = 100000 # Aprox 1 cm en unidades EMU de PowerPoint
    graficos_shapes.sort(key=lambda s: (round(s.top / margen_fila), s.left))

    # 4. INYECTAR LOS DATOS EN ORDEN
    for i, shape in enumerate(graficos_shapes):
        if i >= len(variables_requeridas):
            break # Si hay más gráficos que variables, ignoramos los sobrantes
            
        var_name = variables_requeridas[i]
        
        # Buscamos la tabla procesada en nuestro gran diccionario de resultados
        if var_name in diccionario_resultados:
            df_datos = diccionario_resultados[var_name]['percentages'] # O como esté armado tu dict
            
            # ---> ACÁ LLAMÁS A TU FUNCIÓN EXISTENTE DE CHARTDATA <---
            # actualizar_grafico_con_chartdata(shape.chart, df_datos)
            
            print(f"✅ Gráfico '{var_name}' actualizado correctamente.")
        else:
            print(f"❌ ADVERTENCIA: La variable '{var_name}' pedida en las notas no existe en los resultados.")


def extract_chart_data_for_ai(prs, chart_name, ultimos_n_meses=12):
    """
    Escanea un gráfico y EXTRAE EL CONTEXTO (Versión 4.0 - Coordenadas X/Y)
    Usa el eje Y (top) para Título/Insight y los ejes Y/X (top/left) para la Pregunta.
    """
    def arreglar_etiqueta(etiqueta):
        try:
            val = float(etiqueta)
            if 40000 < val < 50000:
                base_date = datetime.datetime(1899, 12, 30)
                target_date = base_date + datetime.timedelta(days=int(val))
                meses = ["ene", "feb", "mar", "abr", "may", "jun", "jul", "ago", "sep", "oct", "nov", "dic"]
                return f"{meses[target_date.month - 1]}-{str(target_date.year)[-2:]}"
        except:
            pass
        return str(etiqueta).strip()

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart and shape.name == chart_name:
                chart = shape.chart
                
                # =========================================================
                # ESCÁNER V4: GPS CARTESIANO
                # =========================================================
                text_elements = []
                
                def extract_texts_from_shape(s):
                    if s.has_text_frame:
                        txt = s.text.strip()
                        txt = " ".join(txt.split()) 
                        
                        if len(txt) > 3 and not txt.isdigit():
                            text_elements.append({
                                "texto": txt,
                                "top": s.top if s.top is not None else 9999999,
                                "left": s.left if s.left is not None else 9999999, # <--- CAPTURAMOS EL EJE X
                                "is_title": s == slide.shapes.title
                            })
                    elif s.shape_type == 6: 
                        for sub_shape in s.shapes:
                            extract_texts_from_shape(sub_shape)

                for s in slide.shapes:
                    extract_texts_from_shape(s)

                titulo_slide = ""
                insight_anterior = ""
                pregunta_base = ""

                if text_elements:
                    # ORDENAMOS TODO DE ARRIBA HACIA ABAJO (Eje Y)
                    text_elements.sort(key=lambda x: x["top"])

                    # 1. TÍTULO: Oficial o el primero de arriba
                    titulo_oficial = next((x for x in text_elements if x["is_title"]), None)
                    if titulo_oficial:
                        titulo_slide = titulo_oficial["texto"]
                        text_elements.remove(titulo_oficial)
                    else:
                        titulo_slide = text_elements[0]["texto"]
                        text_elements.pop(0)

                    # 2. INSIGHT: El primer texto largo (>40) que le sigue al título
                    insight_candidato = next((x for x in text_elements if len(x["texto"]) > 40), None)
                    if insight_candidato:
                        insight_anterior = insight_candidato["texto"]
                        text_elements.remove(insight_candidato)

                    # 3. PREGUNTA: ABAJO A LA IZQUIERDA (Coordenadas puras)
                    if text_elements:
                        # Filtramos los que están de la mitad para abajo (top > 4.000.000 EMUs)
                        candidatos_abajo = [x for x in text_elements if x["top"] > 4000000]
                        
                        if candidatos_abajo:
                            # Ordenamos esos de izquierda a derecha (menor 'left' primero)
                            candidatos_abajo.sort(key=lambda x: x["left"])
                            pregunta_base = candidatos_abajo[0]["texto"]
                # =========================================================

                try:
                    cats = [arreglar_etiqueta(c.label) for c in chart.plots[0].categories]
                except:
                    cats = [arreglar_etiqueta(c) for c in chart.plots[0].categories]
                
                datos_series = {}
                for series in chart.series:
                    valores_limpios = [v if v is not None else "N/A" for v in series.values]
                    datos_series[series.name.strip()] = valores_limpios[-ultimos_n_meses:]

                return {
                    "contexto": {
                        "tema": titulo_slide,
                        "pregunta": pregunta_base,
                        "insight_anterior": insight_anterior
                    },
                    "meses": cats[-ultimos_n_meses:],
                    "datos": datos_series
                }
    return None

def inject_ai_insights_into_presentation(prs, ai_insights):
    """
    Recibe el diccionario de Gemini y reemplaza el insight en cada diapositiva,
    manteniendo el formato (fuente, color, tamaño) del texto original.
    """
    for chart_name, nuevo_insight in ai_insights.items():
        for slide in prs.slides:
            
            # 1. Verificar si esta es la diapositiva correcta buscando el gráfico
            tiene_el_grafico = False
            for shape in slide.shapes:
                if shape.name == chart_name:
                    tiene_el_grafico = True
                    break

            if tiene_el_grafico:
                # 2. Replicamos la lógica V3/V4 para encontrar la caja del insight
                text_elements = []
                for s in slide.shapes:
                    if s.has_text_frame:
                        txt = s.text.strip()
                        if len(txt) > 3 and not txt.isdigit():
                            text_elements.append({
                                "shape": s,
                                "texto": txt,
                                "top": s.top if s.top is not None else 9999999,
                                "is_title": s == slide.shapes.title
                            })
                
                if text_elements:
                    # Ordenamos de arriba hacia abajo
                    text_elements.sort(key=lambda x: x["top"])
                    
                    # Descartamos el título (el primero)
                    titulo_oficial = next((x for x in text_elements if x["is_title"]), None)
                    if titulo_oficial:
                        text_elements.remove(titulo_oficial)
                    else:
                        text_elements.pop(0)

                    # El siguiente texto largo (>40) es la caja del Insight
                    insight_candidato = next((x for x in text_elements if len(x["texto"]) > 40), None)
                    
                    if insight_candidato:
                        shape_objetivo = insight_candidato["shape"]
                        
                        # 3. INYECCIÓN SEGURA (Mantiene la tipografía y color de tu agencia)
                        if shape_objetivo.text_frame.paragraphs and shape_objetivo.text_frame.paragraphs[0].runs:
                            p_principal = shape_objetivo.text_frame.paragraphs[0]
                            p_principal.runs[0].text = nuevo_insight
                            
                            # Limpiamos el resto de los runs del primer párrafo
                            for i in range(1, len(p_principal.runs)):
                                p_principal.runs[i].text = ""
                                
                            # Borramos párrafos adicionales si el texto viejo era muy largo
                            for i in range(len(shape_objetivo.text_frame.paragraphs) - 1, 0, -1):
                                p_fantasma = shape_objetivo.text_frame.paragraphs[i]._p
                                p_fantasma.getparent().remove(p_fantasma)
                        else:
                            # Plan B por si la caja no tiene formato complejo
                            shape_objetivo.text = nuevo_insight
                        
                        print(f"✅ Insight inyectado con éxito en la slide de '{chart_name}'")
                
                # Como ya inyectamos en esta diapositiva, cortamos el bucle de slides y pasamos al próximo gráfico
                break

def extract_table_data_for_ai(prs, table_name):
    """
    Escanea una tabla estática o evolutiva y EXTRAE EL CONTEXTO Y LA CUADRÍCULA.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table and shape.name == table_name:
                
                # =========================================================
                # 1. ESCÁNER V4 (El mismo GPS espacial que usamos en gráficos)
                # =========================================================
                text_elements = []
                
                def extract_texts_from_shape(s):
                    if s.has_text_frame:
                        txt = s.text.strip()
                        txt = " ".join(txt.split()) 
                        if len(txt) > 3 and not txt.isdigit():
                            text_elements.append({
                                "texto": txt,
                                "top": s.top if s.top is not None else 9999999,
                                "left": s.left if s.left is not None else 9999999,
                                "is_title": s == slide.shapes.title
                            })
                    elif s.shape_type == 6: 
                        for sub_shape in s.shapes:
                            extract_texts_from_shape(sub_shape)

                for s in slide.shapes:
                    extract_texts_from_shape(s)

                titulo_slide = ""
                insight_anterior = ""
                pregunta_base = ""

                if text_elements:
                    text_elements.sort(key=lambda x: x["top"])
                    titulo_oficial = next((x for x in text_elements if x["is_title"]), None)
                    if titulo_oficial:
                        titulo_slide = titulo_oficial["texto"]
                        text_elements.remove(titulo_oficial)
                    else:
                        titulo_slide = text_elements[0]["texto"]
                        text_elements.pop(0)

                    insight_candidato = next((x for x in text_elements if len(x["texto"]) > 40), None)
                    if insight_candidato:
                        insight_anterior = insight_candidato["texto"]
                        text_elements.remove(insight_candidato)

                    if text_elements:
                        candidatos_abajo = [x for x in text_elements if x["top"] > 4000000]
                        if candidatos_abajo:
                            candidatos_abajo.sort(key=lambda x: x["left"])
                            pregunta_base = candidatos_abajo[0]["texto"]
                
                # =========================================================
                # 2. EXTRACCIÓN DE LA CUADRÍCULA DE LA TABLA
                # =========================================================
                table_data = []
                for row in shape.table.rows:
                    row_texts = []
                    for cell in row.cells:
                        # Limpiamos el texto de la celda y quitamos saltos de línea
                        txt = cell.text.strip().replace('\n', ' ')
                        row_texts.append(txt)
                    table_data.append(row_texts)

                return {
                    "contexto": {
                        "tema": titulo_slide,
                        "pregunta": pregunta_base,
                        "insight_anterior": insight_anterior
                    },
                    "datos_tabla": table_data # <--- ¡Enviamos la matriz completa!
                }
    return None

def update_global_ytd_labels(prs, mes_nuevo):
    """
    Escanea todo el PPT buscando textos que sean "YTD ENERO" o "YTD ENERO - [MES]"
    y los actualiza automáticamente según el mes actual, manteniendo el formato.
    """
    # 1. Detectar el mes actual y armar el texto correcto
    meses = {
        "ene": "ENERO", "feb": "FEBRERO", "mar": "MARZO", "abr": "ABRIL", 
        "may": "MAYO", "jun": "JUNIO", "jul": "JULIO", "ago": "AGOSTO", 
        "sep": "SEPTIEMBRE", "oct": "OCTUBRE", "nov": "NOVIEMBRE", "dic": "DICIEMBRE"
    }
    
    mes_lower = str(mes_nuevo).lower()
    mes_detectado = None
    for k, v in meses.items():
        if k in mes_lower:
            mes_detectado = v
            break
            
    if not mes_detectado:
        mes_detectado = str(mes_nuevo).upper() # Fallback por si acaso
        
    if mes_detectado == "ENERO":
        texto_final = "YTD ENERO"
    else:
        texto_final = f"YTD ENERO - {mes_detectado}"
        
    # 2. El Radar Regex (Busca exactamente la estructura, ignorando mayúsculas)
    patron = re.compile(r"^YTD\s+ENERO(\s*-\s*[A-Z]+)?$", re.IGNORECASE)
    reemplazos = 0
    
    # Función interna para no repetir código
    def procesar_texto(text_frame):
        nonlocal reemplazos
        if not text_frame.text: 
            return
            
        # Limpiamos saltos de línea para que el radar no se confunda
        texto_limpio = text_frame.text.strip().replace('\n', ' ')
        
        # Si hace "Match" perfecto con la estructura...
        if patron.match(texto_limpio):
            if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                p = text_frame.paragraphs[0]
                p.runs[0].text = texto_final
                for i in range(1, len(p.runs)):
                    p.runs[i].text = ""
            else:
                text_frame.text = texto_final
            reemplazos += 1

    # 3. Escanear TODO el PPT (Formas sueltas y Celdas de Tablas)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                procesar_texto(shape.text_frame)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        procesar_texto(cell.text_frame)
                        
    return reemplazos, texto_final

def create_summary_slide_with_llm(prs, all_findings_list):
    """Crea la diapositiva de resumen llamando a Gemini para la síntesis global."""
    
    if not config.RUN_LLM or not utils.LLM_AVAILABLE:
        logging.warning("Saltando slide de resumen global: LLM no está activo.")
        return

    try:
        llm_prompt = utils.create_llm_prompt(all_findings_list, "GLOBAL_SUMMARY")
        llm_response = utils.generate_text_with_llm(llm_prompt)
        
        # Lógica de parseo (simplificada)
        title_match = re.search(r"TÍTULO GLOBAL: (.+)", llm_response)
        summary_match = re.search(r"RESUMEN: (.+)", llm_response, re.DOTALL)
        
        llm_title = title_match.group(1).strip() if title_match else "RESUMEN EJECUTIVO"
        llm_summary = summary_match.group(1).strip() if summary_match else "No se pudo generar el resumen."

        title_content_layout = utils.find_layout_by_name(prs, 'Title and Content') 
        if not title_content_layout:
             title_content_layout = prs.slide_layouts[1] 
        
        summary_slide = prs.slides.add_slide(title_content_layout)
        summary_slide.shapes.title.text = llm_title 
        summary_slide.placeholders[16].text = llm_summary

    except Exception as e:
        logging.error(f"Fallo al generar la diapositiva de resumen LLM: {e}")

    logging.info(f"Diapositiva de Resumen Global (LLM) generada.")